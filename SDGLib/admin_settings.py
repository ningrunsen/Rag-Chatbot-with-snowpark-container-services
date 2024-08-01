import streamlit as st
import pandas as pd
import io
import PyPDF2
from docx import Document
from pptx import Presentation
from PIL import Image  # Import PIL for image processing
# import snowflake.connector
# from snowflake.connector import connect
from snowflake.snowpark import Session
import os
from st_snowauth.st_snowauth import snowauth_session


# Create Snowflake session
@st.cache_resource(ttl=600)  # Cache the connection for 10 minutes
def create_snowflake_session():
    session, _ = snowauth_session()
    return session

session = create_snowflake_session()

####################################################################
#########       Confluence: Extractor and Uploader         #########
####################################################################

# Loading Data and Vectorizing will be in SDGLib folder 
# retrieve and display DataFrame from Snowflake
def get_data_from_snowflake():
    SNOWFLAKE_EXTRACTED_DATASET = st.secrets["adminPageSettings"]["snowflake"]["confluence_embed_dataset"]    
    query = f"SELECT * FROM {SNOWFLAKE_EXTRACTED_DATASET}"
    # df = pd.read_sql(query, conn)
    df = session.sql(query).to_pandas()
    # conn.close()
    return df

####################################################################
#########     Staging Documents: Extractor and Uploader    #########
####################################################################
# Theres a weird lag -- I really cant figure it out. Its fine but bothers me
# Its defienty because of the way I wrote this; removing everything after 52 didnt have lag
def stage_documents_handler():
    # Implement User based inputs
    col1, col2, col3 = st.columns(3)
    with col1:
        SNOWFLAKE_DATABASE = st.text_input("Snowflake Database", value=st.secrets["adminPageSettings"]["snowflake"]["default_database"])
    with col2:
        SNOWFLAKE_SCHEMA = st.text_input("Snowflake Schema", value=st.secrets["adminPageSettings"]["snowflake"]["default_schema"])
    with col3:
        SNOWFLAKE_STAGE = st.text_input("Snowflake Stage", value=st.secrets["adminPageSettings"]["snowflake"]["default_stage"])
    # Calling the Dataset and assigning to var
    SNOWFLAKE_EXTRACTED_DATASET = st.secrets["adminPageSettings"]["snowflake"]["extracted_stage_files"]    

    # # Gets Database and Schema from Secrets, can change
    # conn.cursor().execute(f"USE DATABASE {SNOWFLAKE_DATABASE}")
    # conn.cursor().execute(f"USE SCHEMA {SNOWFLAKE_SCHEMA}")
   
    # Ensure the tmp directory exists
    os.makedirs('tmp', exist_ok=True)

    class FileProcessor:
        
        def fetch_file_from_stage(self, full_path: str) -> bytes:
            # Adjusted to use the Snowpark session
            result = session.file.get(f"@{SNOWFLAKE_DATABASE}.{SNOWFLAKE_SCHEMA}.{full_path}", f"file://tmp/")
            local_file_path = f"tmp/{full_path.split('/')[-1]}"
            with open(local_file_path, 'rb') as f:
                return f.read()
            
        # New simplifed logic
        def process_file(self, file_name: str, file_stream: bytes) -> str:
            if file_name.endswith('.xlsx'):
                df = pd.read_excel(io.BytesIO(file_stream))
                return df.to_csv(sep='\t', index=False)
            elif file_name.endswith('.pdf'):
                return "\n".join([page.extract_text() for page in PyPDF2.PdfReader(io.BytesIO(file_stream)).pages])
            elif file_name.endswith('.docx'):
                return "\n".join([para.text for para in Document(io.BytesIO(file_stream)).paragraphs])
            elif file_name.endswith('.csv'):
                df = pd.read_csv(io.BytesIO(file_stream))
                return df.to_csv(sep='\t', index=False)
            elif file_name.endswith('.txt'):
                return io.BytesIO(file_stream).read().decode('utf-8')
            # Figure out how to do images
            elif file_name.endswith(('.jpg', '.png')):
                return f"Image file {file_name} processed."
            elif file_name.endswith('.pptx'):
                return "\n".join([f"Slide {slide_num + 1}:\n" + "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) for slide_num, slide in enumerate(Presentation(io.BytesIO(file_stream)).slides)])
            else:
                return f"Unsupported file type: {file_name}"

        def list_files_in_stage(self):
            # Adjusted to use the Snowpark session
            result = session.sql(f"LIST @{SNOWFLAKE_DATABASE}.{SNOWFLAKE_SCHEMA}.{SNOWFLAKE_STAGE}").collect()
            return result # fetch all details

    processor = FileProcessor()
    # List all files in the Snowflake stage
    with st.spinner("Fetching files from the Snowflake stage..."):
        files_in_stage = processor.list_files_in_stage()
        all_file_urls = [file[0] for file in files_in_stage] # get file paths for the all
        st.info("Select a file from the Snowflake stage or proceed to process all:")
        current_file_url = st.selectbox("Select a file from the Snowflake stage", all_file_urls, label_visibility='collapsed')

    data = [] # Initialize an empty list to collect dat
    col1, col2 = st.columns(2)

    def process_and_store_files(file_urls):
        for file_url in file_urls:
            try:
                file_data = processor.fetch_file_from_stage(file_url)
                file_name = file_url.split('/')[-1]
                selected_file_details = next(file for file in files_in_stage if file[0] == file_url)
                doc_id = selected_file_details[2]
                doc_title = os.path.splitext(file_name)[0]
                source = f'@"{SNOWFLAKE_DATABASE}"."{SNOWFLAKE_SCHEMA}"."{SNOWFLAKE_STAGE}"/{file_url}'
                published = pd.Timestamp(selected_file_details[3])
                doc_content = processor.process_file(file_name, file_data)
                data.append([doc_id, doc_title, doc_content, source, published])
            except Exception as e:
                st.error(f"Error processing file {file_name}: {e}")

        if data:
            columns = ['DOC_ID', 'DOC_TITLE', 'DOC_CONTENT', 'SOURCE', 'PUBLISHED']
            df = pd.DataFrame(data, columns=columns)
            df['PUBLISHED'] = df['PUBLISHED'].dt.strftime('%Y-%m-%d %H:%M:%S')
            # st.write("Document DataFrame:")
            # st.dataframe(df)
            try: # insert df into Snowflake table
                for _, row in df.iterrows():
                    session.sql(f"""
                        INSERT INTO {SNOWFLAKE_EXTRACTED_DATASET} (DOC_ID, DOC_TITLE, DOC_CONTENT, SOURCE, PUBLISHED) 
                        VALUES ('{row['DOC_ID']}', '{row['DOC_TITLE']}', '{row['DOC_CONTENT']}', '{row['SOURCE']}', '{row['PUBLISHED']}')
                    """).collect()
            except Exception as e:
                st.error(f"Error inserting data into Snowflake: {e}")

    # Placeholder for clearing outputs
    output_placeholder = st.empty()

    with col1:
        if st.button("Process File"):
            with st.spinner("Processing and storing file(s)..."):
                process_and_store_files([current_file_url])
    with col2:
        if st.button("Process All Files"):
            with st.spinner("Processing and storing file(s)..."):
                process_and_store_files(all_file_urls)

    if data:
        with output_placeholder.container():
            columns = ['DOC_ID', 'DOC_TITLE', 'DOC_CONTENT', 'SOURCE', 'PUBLISHED']
            df = pd.DataFrame(data, columns=columns)
            df['PUBLISHED'] = df['PUBLISHED'].dt.strftime('%Y-%m-%d %H:%M:%S')
            st.write("Processed Document Data")
            st.dataframe(df)

        
####################################################################
#########       Manual Upload: Extractor and Uploader      #########
####################################################################
def upload_documents_handler():
    # Calling the Dataset and assigning to var
    SNOWFLAKE_EXTRACTED_DATASET = st.secrets["adminPageSettings"]["snowflake"]["extracted_stage_files"] 
    class FileProcessor:
        @staticmethod
        def process_file(file_name: str, file_stream: bytes) -> str:
            if file_name.endswith('.xlsx'):
                df = pd.read_excel(io.BytesIO(file_stream))
                return df.to_csv(sep='\t', index=False)
            elif file_name.endswith('.pdf'):
                return "\n".join([page.extract_text() for page in PyPDF2.PdfReader(io.BytesIO(file_stream)).pages])
            elif file_name.endswith('.docx'):
                return "\n".join([para.text for para in Document(io.BytesIO(file_stream)).paragraphs])
            elif file_name.endswith('.csv'):
                df = pd.read_csv(io.BytesIO(file_stream))
                return df.to_csv(sep='\t', index=False)
            elif file_name.endswith('.txt'):
                return io.BytesIO(file_stream).read().decode('utf-8')
            elif file_name.endswith(('.jpg', '.png')):
                return f"Image file {file_name} processed."
            elif file_name.endswith('.pptx'):
                return "\n".join([f"Slide {slide_num + 1}:\n" + "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) for slide_num, slide in enumerate(Presentation(io.BytesIO(file_stream)).slides)])
            else:
                return f"Unsupported file type: {file_name}"

    data = [] # initialize an empty list to store data for DataFrame

    # Streamlit file uploader
    uploaded_files = st.file_uploader("Upload your files", accept_multiple_files=True)

    # Placeholder for clearing outputs
    output_placeholder = st.empty()

    if uploaded_files:
        with st.spinner("Processing and uploading file(s)..."):
            processor = FileProcessor()
            for uploaded_file in uploaded_files:
                file_name = uploaded_file.name
                file_stream = uploaded_file.read()
                # st.write(f"Processing {file_name}")
                
                # Process each file using the FileProcessor
                doc_id = file_name  # Use file_name as a unique ID for simplicity
                doc_title = os.path.splitext(file_name)[0]
                source = "Uploaded via Streamlit"
                published = pd.Timestamp.now()  # Use current timestamp as published time
                doc_content = processor.process_file(file_name, file_stream)

                # Append the processed data to the list
                data.append([doc_id, doc_title, doc_content, source, published])

            # Create DataFrame from the collected data
            if data:
                columns = ['DOC_ID', 'DOC_TITLE', 'DOC_CONTENT', 'SOURCE', 'PUBLISHED']
                df = pd.DataFrame(data, columns=columns)

                # Convert the timestamp to string format for insertion
                df['PUBLISHED'] = df['PUBLISHED'].dt.strftime('%Y-%m-%d %H:%M:%S')

                # Insert DataFrame into Snowflake table
                try:
                    for _, row in df.iterrows():
                        session.sql(f"""
                            INSERT INTO {SNOWFLAKE_EXTRACTED_DATASET} (DOC_ID, DOC_TITLE, DOC_CONTENT, SOURCE, PUBLISHED) 
                            VALUES ('{row['DOC_ID']}', '{row['DOC_TITLE']}', '{row['DOC_CONTENT']}', '{row['SOURCE']}', '{row['PUBLISHED']}')
                        """).collect()
                except Exception as e:
                    st.error(f"Error inserting data into Snowflake: {e}")

    # Display the processed DataFrame in Streamlit
    if data:
        with output_placeholder.container():
            st.write("Processed Document Data:")
            st.dataframe(df)